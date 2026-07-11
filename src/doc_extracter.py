"""
Extractor de texto para PPTX y DOCX.
Complementa a pdf_extracter.py para soportar más formatos.
"""

from pathlib import Path


def extract_pptx(file_bytes: bytes, filename: str) -> str:
    """Extrae texto de un archivo PPTX."""
    from pptx import Presentation
    import io

    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []

    # Detectar título real: el placeholder de título de la primera diapositiva
    titulo_detectado = None
    if len(prs.slides) > 0:
        primera_slide = prs.slides[0]
        if primera_slide.shapes.title and primera_slide.shapes.title.text.strip():
            titulo_detectado = primera_slide.shapes.title.text.strip()

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs if run.text.strip())
                if line.strip():
                    texts.append(line.strip())

        if texts:
            slides_text.append(f"## Diapositiva {i}\n\n" + "\n".join(texts))

    titulo_final = titulo_detectado or Path(filename).stem
    resultado = f"# {titulo_final}\n\n" + "\n\n".join(slides_text)
    print(f"[PPTX] ✅ Extraídas {len(prs.slides)} diapositivas de {filename}")
    return resultado


def extract_docx(file_bytes: bytes, filename: str) -> str:
    """Extrae texto de un archivo DOCX."""
    from docx import Document
    import io

    doc = Document(io.BytesIO(file_bytes))
    parrafos = []
    titulo_detectado = None

    for para in doc.paragraphs:
        texto = para.text.strip()
        if not texto:
            continue

        # Detectar el título real: estilo "Title" o el primer "Heading 1"
        if titulo_detectado is None and (
            para.style.name.startswith("Title") or para.style.name.startswith("Heading 1")
        ):
            titulo_detectado = texto
            continue  # no lo dupliquemos también como encabezado en el cuerpo

        if para.style.name.startswith("Heading 1"):
            parrafos.append(f"# {texto}")
        elif para.style.name.startswith("Heading 2"):
            parrafos.append(f"## {texto}")
        elif para.style.name.startswith("Heading 3"):
            parrafos.append(f"### {texto}")
        else:
            parrafos.append(texto)

    titulo_final = titulo_detectado or Path(filename).stem
    resultado = f"# {titulo_final}\n\n" + "\n\n".join(parrafos)
    print(f"[DOCX] ✅ Extraídos {len(doc.paragraphs)} párrafos de {filename}")
    return resultado


def extract_file(file_bytes: bytes, filename: str) -> str:
    """
    Router principal: detecta el tipo de archivo y llama al extractor correcto.
    Para PDF usa MinerU API (pdf_extracter.py).
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pptx":
        return extract_pptx(file_bytes, filename)
    elif ext == ".docx":
        return extract_docx(file_bytes, filename)
    elif ext == ".pdf":
        from src.pdf_extracter import extract_pdf_from_bytes
        return extract_pdf_from_bytes(file_bytes, filename)
    else:
        raise ValueError(f"Formato no soportado: {ext}. Usa PDF, PPTX o DOCX.")
